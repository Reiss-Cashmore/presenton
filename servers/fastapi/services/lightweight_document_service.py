"""
Lightweight document converter for Windows/MSIX compatibility.
Uses pure-Python libraries: PyMuPDF (fitz) for PDF, python-docx for DOCX, python-pptx for PPTX.
No subprocess, no external runtimes, MSIX/Appx safe.
"""
import os
from typing import Optional
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation


class LightweightDocumentConverter:
    """Lightweight document converter supporting PDF, DOCX, and PPTX."""
    
    def convert(self, file_path: str) -> str:
        """
        Convert document to markdown text.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text in markdown format
            
        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file does not exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self._convert_pdf(file_path)
        elif file_ext == '.docx':
            return self._convert_docx(file_path)
        elif file_ext == '.pptx':
            return self._convert_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _convert_pdf(self, path: str) -> str:
        """
        Convert PDF to markdown using PyMuPDF.
        
        Args:
            path: Path to PDF file
            
        Returns:
            Extracted text in markdown format
        """
        doc = fitz.open(path)
        markdown_parts = []
        
        for page_num, page in enumerate(doc, start=1):
            # Extract text blocks for better structure preservation
            blocks = page.get_text("blocks")
            
            page_text_parts = []
            for block in blocks:
                # block format: (x0, y0, x1, y1, "text", block_type, block_no)
                if len(block) >= 5:
                    text = block[4].strip()
                    if text:
                        page_text_parts.append(text)
            
            if page_text_parts:
                page_content = "\n\n".join(page_text_parts)
                # Add page separator for multi-page documents
                if len(doc) > 1:
                    markdown_parts.append(f"## Page {page_num}\n\n{page_content}")
                else:
                    markdown_parts.append(page_content)
        
        doc.close()
        return "\n\n".join(markdown_parts)
    
    def _convert_docx(self, path: str) -> str:
        """
        Convert DOCX to markdown using python-docx.
        
        Args:
            path: Path to DOCX file
            
        Returns:
            Extracted text in markdown format
        """
        doc = DocxDocument(path)
        markdown_parts = []
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            
            # Check if it's a heading based on style
            style_name = paragraph.style.name.lower()
            
            if 'heading' in style_name:
                # Extract heading level from style (e.g., "Heading 1" -> 1)
                try:
                    level = int(style_name.split()[-1])
                    level = min(level, 6)  # Markdown supports up to 6 levels
                    markdown_parts.append(f"{'#' * level} {text}")
                except (ValueError, IndexError):
                    # Fallback: treat as regular paragraph
                    markdown_parts.append(text)
            else:
                markdown_parts.append(text)
        
        return "\n\n".join(markdown_parts)
    
    def _convert_pptx(self, path: str) -> str:
        """
        Convert PPTX to markdown using python-pptx.
        
        Args:
            path: Path to PPTX file
            
        Returns:
            Extracted text in markdown format
        """
        prs = Presentation(path)
        markdown_parts = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            
            # Extract slide title (usually first shape with title placeholder)
            title_text = None
            for shape in slide.shapes:
                if hasattr(shape, "placeholder"):
                    if shape.placeholder.placeholder_format.type == 1:  # Title placeholder
                        if hasattr(shape, "text") and shape.text.strip():
                            title_text = shape.text.strip()
                            break
            
            # If no title placeholder found, try to find text box at top
            if not title_text:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's likely a title (first text shape, short text)
                        text = shape.text.strip()
                        if len(text) < 200:  # Heuristic: titles are usually short
                            title_text = text
                            break
            
            # Add slide title
            if title_text:
                slide_parts.append(f"# {title_text}")
            else:
                slide_parts.append(f"# Slide {slide_num}")
            
            # Extract content (bullet points and text)
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                
                text = shape.text.strip()
                if not text:
                    continue
                
                # Skip if this is the title we already added
                if title_text and text == title_text:
                    continue
                
                # Check if it's a text frame with paragraphs (bullet points)
                if hasattr(shape, "text_frame"):
                    paragraphs = shape.text_frame.paragraphs
                    if len(paragraphs) > 1:
                        # Multiple paragraphs - likely bullet points
                        for para in paragraphs:
                            para_text = para.text.strip()
                            if para_text:
                                # Check bullet level
                                level = para.level
                                indent = "  " * level
                                slide_parts.append(f"{indent}- {para_text}")
                    else:
                        # Single paragraph
                        if text and text != title_text:
                            slide_parts.append(text)
                else:
                    # Plain text shape
                    if text and text != title_text:
                        slide_parts.append(text)
            
            if slide_parts:
                markdown_parts.append("\n".join(slide_parts))
        
        return "\n\n---\n\n".join(markdown_parts)


class DocumentService:
    """
    Document service wrapper providing parse_to_markdown interface.
    Compatible with DoclingService interface for easy swapping.
    """
    
    def __init__(self):
        self.converter = LightweightDocumentConverter()
    
    def parse_to_markdown(self, file_path: str) -> str:
        """
        Parse document to markdown format.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text in markdown format
        """
        return self.converter.convert(file_path)
